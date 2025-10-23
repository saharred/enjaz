"""
Enhanced module for generating PowerPoint presentations with charts.
Uses python-pptx library with Qatar visual identity and matplotlib for charts.
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from enjaz.analysis import get_band
import matplotlib.pyplot as plt
import matplotlib
matplotlib.use('Agg')  # Non-interactive backend
import pandas as pd
import tempfile
import os
from enjaz.arabic_text_helper import fix_arabic_text

# Qatar Colors
MAROON = RGBColor(138, 21, 56)  # #8A1538
GOLD = RGBColor(201, 162, 39)   # #C9A227
WHITE = RGBColor(255, 255, 255)
DARK_GRAY = RGBColor(51, 51, 51)
LIGHT_GRAY = RGBColor(245, 245, 245)

# Colors for charts (hex)
MAROON_HEX = '#8A1538'
GOLD_HEX = '#C9A227'


def create_band_distribution_chart(band_distribution, output_path):
    """Create a pie chart for band distribution."""
    # Configure matplotlib for Arabic
    plt.rcParams['font.family'] = 'DejaVu Sans'
    
    # Prepare data
    labels = []
    sizes = []
    colors = ['#FFD700', '#C9A227', '#CD7F32', '#8B4513', '#FF6347', '#DC143C']
    
    band_names = [
        'البلاتينية',
        'الذهبية', 
        'الفضية',
        'البرونزية',
        'يحتاج إلى تطوير',
        'لا يستفيد من النظام'
    ]
    
    for idx, band in enumerate(band_names):
        count = band_distribution.get(band, 0)
        if band == 'يحتاج إلى تطوير':
            count = band_distribution.get('يحتاج إلى تطوير من النظام', count)
        if count > 0:
            labels.append(fix_arabic_text(band))
            sizes.append(count)
    
    # Create figure
    fig, ax = plt.subplots(figsize=(10, 7))
    
    # Create pie chart
    wedges, texts, autotexts = ax.pie(
        sizes,
        labels=labels,
        autopct='%1.1f%%',
        startangle=90,
        colors=colors[:len(sizes)],
        textprops={'fontsize': 14, 'weight': 'bold'}
    )
    
    # Enhance text
    for autotext in autotexts:
        autotext.set_color('white')
        autotext.set_fontsize(16)
        autotext.set_weight('bold')
    
    ax.axis('equal')
    plt.title(fix_arabic_text('توزيع الطلاب حسب الفئات'), fontsize=20, weight='bold', pad=20)
    plt.tight_layout()
    
    # Save
    plt.savefig(output_path, dpi=150, bbox_inches='tight', facecolor='white')
    plt.close()
    
    return output_path


def create_subject_comparison_chart(subject_stats, output_path):
    """Create a bar chart for subject comparison."""
    if not subject_stats:
        return None
    
    # Prepare data - top 6 subjects
    subjects = [fix_arabic_text(s['subject_name']) for s in subject_stats[:6]]
    rates = [s['completion_rate'] for s in subject_stats[:6]]
    
    # Create figure
    fig, ax = plt.subplots(figsize=(10, 6))
    
    # Create bars with gradient colors
    bars = ax.barh(subjects, rates, color=MAROON_HEX, edgecolor=GOLD_HEX, linewidth=2)
    
    # Add value labels
    for i, (bar, rate) in enumerate(zip(bars, rates)):
        ax.text(rate + 1, i, f'{rate:.1f}%', 
                va='center', fontsize=12, weight='bold', color=MAROON_HEX)
    
    # Styling
    ax.set_xlabel(fix_arabic_text('نسبة الإنجاز (%)'), fontsize=14, weight='bold')
    ax.set_title(fix_arabic_text('مقارنة أداء المواد الدراسية'), fontsize=18, weight='bold', pad=15)
    ax.set_xlim(0, 105)
    ax.grid(axis='x', alpha=0.3, linestyle='--')
    ax.spines['top'].set_visible(False)
    ax.spines['right'].set_visible(False)
    
    plt.tight_layout()
    plt.savefig(output_path, dpi=150, bbox_inches='tight', facecolor='white')
    plt.close()
    
    return output_path


def create_statistics_chart(school_stats, output_path):
    """Create a visual chart for key statistics."""
    fig, ax = plt.subplots(figsize=(10, 6))
    
    # Data
    categories = [fix_arabic_text('إجمالي الطلاب'), fix_arabic_text('إجمالي التقييمات'), fix_arabic_text('التقييمات المنجزة')]
    values = [
        school_stats['total_students'],
        school_stats['total_assessments'],
        school_stats['total_completed']
    ]
    
    # Create bars
    bars = ax.bar(categories, values, color=[MAROON_HEX, GOLD_HEX, '#6B8E23'], 
                   edgecolor='black', linewidth=1.5, alpha=0.8)
    
    # Add value labels on top of bars
    for bar, value in zip(bars, values):
        height = bar.get_height()
        ax.text(bar.get_x() + bar.get_width()/2., height,
                f'{int(value)}',
                ha='center', va='bottom', fontsize=16, weight='bold')
    
    # Styling
    ax.set_ylabel(fix_arabic_text('العدد'), fontsize=14, weight='bold')
    ax.set_title(fix_arabic_text('الإحصائيات الرئيسية'), fontsize=18, weight='bold', pad=15)
    ax.spines['top'].set_visible(False)
    ax.spines['right'].set_visible(False)
    ax.grid(axis='y', alpha=0.3, linestyle='--')
    
    plt.tight_layout()
    plt.savefig(output_path, dpi=150, bbox_inches='tight', facecolor='white')
    plt.close()
    
    return output_path


def create_title_slide(prs, title, subtitle):
    """Create title slide with Qatar visual identity."""
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)
    
    # Background
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = MAROON
    
    # Title
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(2.5),
        Inches(9), Inches(1.5)
    )
    title_frame = title_box.text_frame
    title_frame.text = title
    title_para = title_frame.paragraphs[0]
    title_para.alignment = PP_ALIGN.CENTER
    title_para.font.size = Pt(54)
    title_para.font.bold = True
    title_para.font.color.rgb = GOLD
    title_para.font.name = 'Cairo'
    
    # Subtitle
    subtitle_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(4.2),
        Inches(9), Inches(1)
    )
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = subtitle
    subtitle_para = subtitle_frame.paragraphs[0]
    subtitle_para.alignment = PP_ALIGN.CENTER
    subtitle_para.font.size = Pt(28)
    subtitle_para.font.color.rgb = WHITE
    subtitle_para.font.name = 'Cairo'
    
    return slide


def create_content_slide(prs, title):
    """Create a content slide with header."""
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)
    
    # Background
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = WHITE
    
    # Header bar
    header = slide.shapes.add_shape(
        1,  # Rectangle
        Inches(0), Inches(0),
        Inches(10), Inches(1)
    )
    header.fill.solid()
    header.fill.fore_color.rgb = MAROON
    header.line.fill.background()
    
    # Title in header
    title_frame = header.text_frame
    title_frame.text = title
    title_para = title_frame.paragraphs[0]
    title_para.alignment = PP_ALIGN.CENTER
    title_para.font.size = Pt(36)
    title_para.font.bold = True
    title_para.font.color.rgb = GOLD
    title_para.font.name = 'Cairo'
    
    return slide


def add_statistics_slide_with_chart(prs, school_stats):
    """Add statistics slide with chart."""
    slide = create_content_slide(prs, "📊 الإحصائيات الرئيسية")
    
    # Create chart
    with tempfile.NamedTemporaryFile(delete=False, suffix='.png') as tmp:
        chart_path = create_statistics_chart(school_stats, tmp.name)
        
        # Add chart to slide
        slide.shapes.add_picture(
            chart_path,
            Inches(0.5), Inches(1.5),
            width=Inches(9)
        )
        
        # Clean up
        os.unlink(chart_path)
    
    # Add completion rate text
    completion_box = slide.shapes.add_textbox(
        Inches(3), Inches(6.5),
        Inches(4), Inches(0.8)
    )
    completion_frame = completion_box.text_frame
    completion_frame.text = f"نسبة الإنجاز الكلية: {school_stats['completion_rate']:.1f}%"
    completion_para = completion_frame.paragraphs[0]
    completion_para.alignment = PP_ALIGN.CENTER
    completion_para.font.size = Pt(24)
    completion_para.font.bold = True
    completion_para.font.color.rgb = MAROON
    completion_para.font.name = 'Cairo'
    
    return slide


def add_band_distribution_slide_with_chart(prs, school_stats):
    """Add band distribution slide with pie chart."""
    slide = create_content_slide(prs, "📈 توزيع الطلاب حسب الفئات")
    
    # Create chart
    with tempfile.NamedTemporaryFile(delete=False, suffix='.png') as tmp:
        chart_path = create_band_distribution_chart(school_stats['band_distribution'], tmp.name)
        
        # Add chart to slide
        slide.shapes.add_picture(
            chart_path,
            Inches(1), Inches(1.3),
            width=Inches(8)
        )
        
        # Clean up
        os.unlink(chart_path)
    
    return slide


def add_subject_analysis_slide_with_chart(prs, subject_stats):
    """Add subject analysis slide with bar chart."""
    slide = create_content_slide(prs, "📚 تحليل المواد الدراسية")
    
    if not subject_stats:
        return slide
    
    # Create chart
    with tempfile.NamedTemporaryFile(delete=False, suffix='.png') as tmp:
        chart_path = create_subject_comparison_chart(subject_stats, tmp.name)
        
        if chart_path:
            # Add chart to slide
            slide.shapes.add_picture(
                chart_path,
                Inches(0.5), Inches(1.3),
                width=Inches(9)
            )
            
            # Clean up
            os.unlink(chart_path)
    
    return slide


def add_recommendations_slide(prs, completion_rate):
    """Add recommendations slide."""
    slide = create_content_slide(prs, "💡 التوصيات التلقائية")
    
    # Intro text
    intro_box = slide.shapes.add_textbox(
        Inches(1), Inches(1.3),
        Inches(8), Inches(0.6)
    )
    intro_frame = intro_box.text_frame
    band = get_band(completion_rate)
    intro_frame.text = f"بناءً على نسبة الإنجاز الكلية ({completion_rate:.1f}% - {band}):"
    intro_para = intro_frame.paragraphs[0]
    intro_para.alignment = PP_ALIGN.CENTER
    intro_para.font.size = Pt(22)
    intro_para.font.bold = True
    intro_para.font.color.rgb = MAROON
    intro_para.font.name = 'Cairo'
    
    # Recommendations
    recommendations = [
        "الحفاظ على الأداء الجيد والعمل على تعزيز ثقافة الإنجاز",
        "تحفيز الطلاب المتميزين وتكريمهم لتشجيع الآخرين",
        "توفير برامج دعم إضافية للطلاب المتعثرين",
        "تفعيل التواصل مع أولياء الأمور وإرسال تقارير دورية",
        "المتابعة الدورية الأسبوعية لقياس التحسن"
    ]
    
    y_start = Inches(2.2)
    for idx, rec in enumerate(recommendations):
        # Number circle
        circle = slide.shapes.add_shape(
            9,  # Oval
            Inches(0.8), y_start + Inches(idx * 0.85),
            Inches(0.5), Inches(0.5)
        )
        circle.fill.solid()
        circle.fill.fore_color.rgb = MAROON
        circle.line.fill.background()
        
        # Number text
        num_frame = circle.text_frame
        num_frame.text = str(idx + 1)
        num_para = num_frame.paragraphs[0]
        num_para.alignment = PP_ALIGN.CENTER
        num_para.font.size = Pt(20)
        num_para.font.bold = True
        num_para.font.color.rgb = GOLD
        num_para.font.name = 'Cairo'
        
        # Recommendation text
        rec_box = slide.shapes.add_textbox(
            Inches(1.5), y_start + Inches(idx * 0.85),
            Inches(7.5), Inches(0.7)
        )
        rec_frame = rec_box.text_frame
        rec_frame.text = rec
        rec_frame.word_wrap = True
        rec_para = rec_frame.paragraphs[0]
        rec_para.alignment = PP_ALIGN.RIGHT
        rec_para.font.size = Pt(18)
        rec_para.font.color.rgb = DARK_GRAY
        rec_para.font.name = 'Cairo'
    
    return slide


def add_coordinator_actions_slide(prs, actions_text):
    """Add coordinator actions slide."""
    slide = create_content_slide(prs, "📝 إجراءات منسق المشاريع")
    
    # Actions text box
    actions_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(1.3),
        Inches(9), Inches(5.5)
    )
    actions_frame = actions_box.text_frame
    actions_frame.text = actions_text
    actions_frame.word_wrap = True
    
    for para in actions_frame.paragraphs:
        para.alignment = PP_ALIGN.RIGHT
        para.font.size = Pt(16)
        para.font.color.rgb = DARK_GRAY
        para.font.name = 'Cairo'
        para.line_spacing = 1.5
    
    return slide


def add_closing_slide(prs):
    """Add closing slide."""
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)
    
    # Background
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = MAROON
    
    # Thank you message
    thanks_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(2.5),
        Inches(9), Inches(1.5)
    )
    thanks_frame = thanks_box.text_frame
    thanks_frame.text = "شكراً لكم"
    thanks_para = thanks_frame.paragraphs[0]
    thanks_para.alignment = PP_ALIGN.CENTER
    thanks_para.font.size = Pt(60)
    thanks_para.font.bold = True
    thanks_para.font.color.rgb = GOLD
    thanks_para.font.name = 'Cairo'
    
    # Contact info
    contact_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(4.5),
        Inches(9), Inches(1.5)
    )
    contact_frame = contact_box.text_frame
    contact_text = "تطوير وتنفيذ: Sahar Osman\nE-Learning Projects Coordinator\ns.mahgoub0101@education.qa\nlinkedin.com/in/sahar-osman-a19a45209"
    contact_frame.text = contact_text
    
    for para in contact_frame.paragraphs:
        para.alignment = PP_ALIGN.CENTER
        para.font.size = Pt(18)
        para.font.color.rgb = WHITE
        para.font.name = 'Cairo'
    
    return slide


def generate_school_presentation(school_stats, coordinator_actions="", subject_stats=None, output_path="school_report.pptx"):
    """
    Generate complete school presentation with charts.
    
    Args:
        school_stats: Dictionary with school statistics
        coordinator_actions: Text with coordinator actions
        subject_stats: List of subject statistics dictionaries
        output_path: Path to save the presentation
    
    Returns:
        Path to the generated presentation
    """
    # Create presentation
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # Add slides
    create_title_slide(
        prs,
        "تقرير المدرسة الشامل",
        "مدرسة عثمان بن عفّان النموذجية للبنين\nوزارة التعليم والتعليم العالي"
    )
    
    add_statistics_slide_with_chart(prs, school_stats)
    add_band_distribution_slide_with_chart(prs, school_stats)
    
    # Add subject analysis if available
    if subject_stats:
        add_subject_analysis_slide_with_chart(prs, subject_stats)
    
    add_recommendations_slide(prs, school_stats['completion_rate'])
    
    if coordinator_actions.strip():
        add_coordinator_actions_slide(prs, coordinator_actions)
    
    add_closing_slide(prs)
    
    # Save presentation
    prs.save(output_path)
    
    return output_path

