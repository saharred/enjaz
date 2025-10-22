"""
Module for generating PowerPoint presentations for school reports.
Uses python-pptx library with Qatar visual identity.
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from enjaz.analysis import get_band
import pandas as pd


# Qatar Colors
MAROON = RGBColor(138, 21, 56)  # #8A1538
GOLD = RGBColor(201, 162, 39)   # #C9A227
WHITE = RGBColor(255, 255, 255)
DARK_GRAY = RGBColor(51, 51, 51)
LIGHT_GRAY = RGBColor(245, 245, 245)


def create_title_slide(prs, title, subtitle):
    """Create title slide with Qatar visual identity."""
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)
    
    # Background
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = MAROON
    
    # Title
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(2.5),
        Inches(9), Inches(1.5)
    )
    title_frame = title_box.text_frame
    title_frame.text = title
    title_para = title_frame.paragraphs[0]
    title_para.alignment = PP_ALIGN.CENTER
    title_para.font.size = Pt(54)
    title_para.font.bold = True
    title_para.font.color.rgb = GOLD
    title_para.font.name = 'Cairo'
    
    # Subtitle
    subtitle_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(4.2),
        Inches(9), Inches(1)
    )
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = subtitle
    subtitle_para = subtitle_frame.paragraphs[0]
    subtitle_para.alignment = PP_ALIGN.CENTER
    subtitle_para.font.size = Pt(28)
    subtitle_para.font.color.rgb = WHITE
    subtitle_para.font.name = 'Cairo'
    
    return slide


def create_content_slide(prs, title):
    """Create a content slide with header."""
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)
    
    # Background
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = WHITE
    
    # Header bar
    header = slide.shapes.add_shape(
        1,  # Rectangle
        Inches(0), Inches(0),
        Inches(10), Inches(1)
    )
    header.fill.solid()
    header.fill.fore_color.rgb = MAROON
    header.line.fill.background()
    
    # Title in header
    title_frame = header.text_frame
    title_frame.text = title
    title_para = title_frame.paragraphs[0]
    title_para.alignment = PP_ALIGN.CENTER
    title_para.font.size = Pt(36)
    title_para.font.bold = True
    title_para.font.color.rgb = GOLD
    title_para.font.name = 'Cairo'
    
    return slide


def add_statistics_slide(prs, school_stats):
    """Add key statistics slide."""
    slide = create_content_slide(prs, "📊 الإحصائيات الرئيسية")
    
    # Statistics cards
    stats_data = [
        ("👥 إجمالي الطلاب", str(school_stats['total_students'])),
        ("📊 إجمالي التقييمات", str(school_stats['total_assessments'])),
        ("✅ التقييمات المُنجزة", str(school_stats['total_completed'])),
        ("🎯 نسبة الإنجاز", f"{school_stats['completion_rate']:.1f}%")
    ]
    
    # Create 2x2 grid
    x_positions = [Inches(0.5), Inches(5.5)]
    y_positions = [Inches(1.5), Inches(4)]
    
    idx = 0
    for row in range(2):
        for col in range(2):
            if idx < len(stats_data):
                label, value = stats_data[idx]
                
                # Card background
                card = slide.shapes.add_shape(
                    1,  # Rectangle
                    x_positions[col], y_positions[row],
                    Inches(4), Inches(2)
                )
                card.fill.solid()
                card.fill.fore_color.rgb = LIGHT_GRAY
                card.line.color.rgb = MAROON
                card.line.width = Pt(3)
                
                # Label
                label_box = slide.shapes.add_textbox(
                    x_positions[col] + Inches(0.2), y_positions[row] + Inches(0.3),
                    Inches(3.6), Inches(0.6)
                )
                label_frame = label_box.text_frame
                label_frame.text = label
                label_para = label_frame.paragraphs[0]
                label_para.alignment = PP_ALIGN.CENTER
                label_para.font.size = Pt(20)
                label_para.font.color.rgb = MAROON
                label_para.font.name = 'Cairo'
                
                # Value
                value_box = slide.shapes.add_textbox(
                    x_positions[col] + Inches(0.2), y_positions[row] + Inches(1),
                    Inches(3.6), Inches(0.8)
                )
                value_frame = value_box.text_frame
                value_frame.text = value
                value_para = value_frame.paragraphs[0]
                value_para.alignment = PP_ALIGN.CENTER
                value_para.font.size = Pt(40)
                value_para.font.bold = True
                value_para.font.color.rgb = MAROON
                value_para.font.name = 'Cairo'
                
                idx += 1
    
    return slide


def add_band_distribution_slide(prs, school_stats):
    """Add band distribution slide."""
    slide = create_content_slide(prs, "📈 توزيع الطلاب حسب الفئات")
    
    # Band data
    bands = [
        ("البلاتينية", school_stats['band_distribution'].get('البلاتينية', 0)),
        ("الذهبية", school_stats['band_distribution'].get('الذهبية', 0)),
        ("الفضية", school_stats['band_distribution'].get('الفضية', 0)),
        ("البرونزية", school_stats['band_distribution'].get('البرونزية', 0)),
        ("يحتاج إلى تطوير", school_stats['band_distribution'].get('يحتاج إلى تطوير من النظام', 0)),
        ("لا يستفيد من النظام", school_stats['band_distribution'].get('لا يستفيد من النظام', 0))
    ]
    
    # Create table
    y_start = Inches(1.5)
    total_students = max(school_stats['total_students'], 1)
    
    for idx, (band_name, count) in enumerate(bands):
        percentage = (count / total_students) * 100
        
        # Band name
        name_box = slide.shapes.add_textbox(
            Inches(1), y_start + Inches(idx * 0.8),
            Inches(3), Inches(0.6)
        )
        name_frame = name_box.text_frame
        name_frame.text = band_name
        name_para = name_frame.paragraphs[0]
        name_para.alignment = PP_ALIGN.RIGHT
        name_para.font.size = Pt(24)
        name_para.font.bold = True
        name_para.font.color.rgb = MAROON
        name_para.font.name = 'Cairo'
        
        # Count
        count_box = slide.shapes.add_textbox(
            Inches(4.5), y_start + Inches(idx * 0.8),
            Inches(2), Inches(0.6)
        )
        count_frame = count_box.text_frame
        count_frame.text = str(count)
        count_para = count_frame.paragraphs[0]
        count_para.alignment = PP_ALIGN.CENTER
        count_para.font.size = Pt(28)
        count_para.font.bold = True
        count_para.font.color.rgb = DARK_GRAY
        count_para.font.name = 'Cairo'
        
        # Percentage
        pct_box = slide.shapes.add_textbox(
            Inches(7), y_start + Inches(idx * 0.8),
            Inches(2), Inches(0.6)
        )
        pct_frame = pct_box.text_frame
        pct_frame.text = f"{percentage:.1f}%"
        pct_para = pct_frame.paragraphs[0]
        pct_para.alignment = PP_ALIGN.CENTER
        pct_para.font.size = Pt(28)
        pct_para.font.color.rgb = GOLD
        pct_para.font.name = 'Cairo'
    
    return slide


def add_subject_analysis_slide(prs, subject_stats):
    """Add subject analysis slide."""
    slide = create_content_slide(prs, "📚 تحليل المواد الدراسية")
    
    if not subject_stats:
        return slide
    
    # Show top 5 subjects
    top_subjects = subject_stats[:min(5, len(subject_stats))]
    
    # Create table
    y_start = Inches(1.5)
    
    # Headers
    headers = ["المادة", "الطلاب", "نسبة الإنجاز", "الفئة"]
    x_positions = [Inches(0.5), Inches(3.5), Inches(5.5), Inches(7.5)]
    widths = [Inches(2.8), Inches(1.8), Inches(1.8), Inches(2)]
    
    for idx, (header, x_pos, width) in enumerate(zip(headers, x_positions, widths)):
        header_box = slide.shapes.add_textbox(x_pos, y_start, width, Inches(0.5))
        header_frame = header_box.text_frame
        header_frame.text = header
        header_para = header_frame.paragraphs[0]
        header_para.alignment = PP_ALIGN.CENTER
        header_para.font.size = Pt(20)
        header_para.font.bold = True
        header_para.font.color.rgb = MAROON
        header_para.font.name = 'Cairo'
    
    # Data rows
    y_start += Inches(0.7)
    for idx, subject in enumerate(top_subjects):
        y_pos = y_start + Inches(idx * 0.7)
        
        # Subject name
        name_box = slide.shapes.add_textbox(x_positions[0], y_pos, widths[0], Inches(0.5))
        name_frame = name_box.text_frame
        name_frame.text = subject['subject_name']
        name_para = name_frame.paragraphs[0]
        name_para.alignment = PP_ALIGN.RIGHT
        name_para.font.size = Pt(18)
        name_para.font.color.rgb = DARK_GRAY
        name_para.font.name = 'Cairo'
        
        # Students count
        count_box = slide.shapes.add_textbox(x_positions[1], y_pos, widths[1], Inches(0.5))
        count_frame = count_box.text_frame
        count_frame.text = str(subject['total_students'])
        count_para = count_frame.paragraphs[0]
        count_para.alignment = PP_ALIGN.CENTER
        count_para.font.size = Pt(18)
        count_para.font.color.rgb = DARK_GRAY
        count_para.font.name = 'Cairo'
        
        # Completion rate
        rate_box = slide.shapes.add_textbox(x_positions[2], y_pos, widths[2], Inches(0.5))
        rate_frame = rate_box.text_frame
        rate_frame.text = f"{subject['completion_rate']:.1f}%"
        rate_para = rate_frame.paragraphs[0]
        rate_para.alignment = PP_ALIGN.CENTER
        rate_para.font.size = Pt(20)
        rate_para.font.bold = True
        rate_para.font.color.rgb = GOLD
        rate_para.font.name = 'Cairo'
        
        # Band
        band_box = slide.shapes.add_textbox(x_positions[3], y_pos, widths[3], Inches(0.5))
        band_frame = band_box.text_frame
        band_frame.text = subject['band']
        band_para = band_frame.paragraphs[0]
        band_para.alignment = PP_ALIGN.CENTER
        band_para.font.size = Pt(16)
        band_para.font.color.rgb = MAROON
        band_para.font.name = 'Cairo'
    
    return slide


def add_recommendations_slide(prs, completion_rate):
    """Add recommendations slide."""
    slide = create_content_slide(prs, "💡 التوصيات التلقائية")
    
    # Intro text
    intro_box = slide.shapes.add_textbox(
        Inches(1), Inches(1.3),
        Inches(8), Inches(0.6)
    )
    intro_frame = intro_box.text_frame
    band = get_band(completion_rate)
    intro_frame.text = f"بناءً على نسبة الإنجاز الكلية ({completion_rate:.1f}% - {band}):"
    intro_para = intro_frame.paragraphs[0]
    intro_para.alignment = PP_ALIGN.CENTER
    intro_para.font.size = Pt(22)
    intro_para.font.bold = True
    intro_para.font.color.rgb = MAROON
    intro_para.font.name = 'Cairo'
    
    # Recommendations
    recommendations = [
        "الحفاظ على الأداء الجيد والعمل على تعزيز ثقافة الإنجاز",
        "تحفيز الطلاب المتميزين وتكريمهم لتشجيع الآخرين",
        "توفير برامج دعم إضافية للطلاب المتعثرين",
        "تفعيل التواصل مع أولياء الأمور وإرسال تقارير دورية",
        "المتابعة الدورية الأسبوعية لقياس التحسن"
    ]
    
    y_start = Inches(2.2)
    for idx, rec in enumerate(recommendations):
        # Number circle
        circle = slide.shapes.add_shape(
            9,  # Oval
            Inches(0.8), y_start + Inches(idx * 0.85),
            Inches(0.5), Inches(0.5)
        )
        circle.fill.solid()
        circle.fill.fore_color.rgb = MAROON
        circle.line.fill.background()
        
        # Number text
        num_frame = circle.text_frame
        num_frame.text = str(idx + 1)
        num_para = num_frame.paragraphs[0]
        num_para.alignment = PP_ALIGN.CENTER
        num_para.font.size = Pt(20)
        num_para.font.bold = True
        num_para.font.color.rgb = GOLD
        num_para.font.name = 'Cairo'
        
        # Recommendation text
        rec_box = slide.shapes.add_textbox(
            Inches(1.5), y_start + Inches(idx * 0.85),
            Inches(7.5), Inches(0.7)
        )
        rec_frame = rec_box.text_frame
        rec_frame.text = rec
        rec_frame.word_wrap = True
        rec_para = rec_frame.paragraphs[0]
        rec_para.alignment = PP_ALIGN.RIGHT
        rec_para.font.size = Pt(18)
        rec_para.font.color.rgb = DARK_GRAY
        rec_para.font.name = 'Cairo'
    
    return slide


def add_coordinator_actions_slide(prs, actions_text):
    """Add coordinator actions slide."""
    slide = create_content_slide(prs, "📝 إجراءات منسق المشاريع")
    
    # Actions text box
    actions_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(1.3),
        Inches(9), Inches(5)
    )
    actions_frame = actions_box.text_frame
    actions_frame.text = actions_text
    actions_frame.word_wrap = True
    
    for para in actions_frame.paragraphs:
        para.alignment = PP_ALIGN.RIGHT
        para.font.size = Pt(18)
        para.font.color.rgb = DARK_GRAY
        para.font.name = 'Cairo'
        para.line_spacing = 1.5
    
    return slide


def add_closing_slide(prs):
    """Add closing slide."""
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)
    
    # Background
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = MAROON
    
    # Thank you message
    thanks_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(2.5),
        Inches(9), Inches(1.5)
    )
    thanks_frame = thanks_box.text_frame
    thanks_frame.text = "شكراً لكم"
    thanks_para = thanks_frame.paragraphs[0]
    thanks_para.alignment = PP_ALIGN.CENTER
    thanks_para.font.size = Pt(60)
    thanks_para.font.bold = True
    thanks_para.font.color.rgb = GOLD
    thanks_para.font.name = 'Cairo'
    
    # Contact info
    contact_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(4.5),
        Inches(9), Inches(1.5)
    )
    contact_frame = contact_box.text_frame
    contact_text = "تطوير وتنفيذ: Sahar Osman\nE-Learning Projects Coordinator\ns.mahgoub0101@education.qa"
    contact_frame.text = contact_text
    
    for para in contact_frame.paragraphs:
        para.alignment = PP_ALIGN.CENTER
        para.font.size = Pt(20)
        para.font.color.rgb = WHITE
        para.font.name = 'Cairo'
    
    return slide


def generate_school_presentation(school_stats, coordinator_actions="", subject_stats=None, output_path="school_report.pptx"):
    """
    Generate complete school presentation.
    
    Args:
        school_stats: Dictionary with school statistics
        coordinator_actions: Text with coordinator actions
        output_path: Path to save the presentation
    
    Returns:
        Path to the generated presentation
    """
    # Create presentation
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # Add slides
    create_title_slide(
        prs,
        "تقرير المدرسة الشامل",
        "مدرسة عثمان بن عفّان النموذجية للبنين\nوزارة التعليم والتعليم العالي – دولة قطر"
    )
    
    add_statistics_slide(prs, school_stats)
    add_band_distribution_slide(prs, school_stats)
    
    # Add subject analysis if available
    if subject_stats:
        add_subject_analysis_slide(prs, subject_stats)
    
    add_recommendations_slide(prs, school_stats['completion_rate'])
    
    if coordinator_actions.strip():
        add_coordinator_actions_slide(prs, coordinator_actions)
    
    add_closing_slide(prs)
    
    # Save presentation
    prs.save(output_path)
    
    return output_path

